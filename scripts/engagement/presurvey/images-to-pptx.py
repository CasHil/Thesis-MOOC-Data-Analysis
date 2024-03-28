import os
import argparse
from pptx import Presentation, presentation, shapes
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx.slide import Slide, Slides
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.text.text import TextFrame


def add_title_slide(centered_image_presentation: presentation.Presentation, title_text: str) -> None:

    slide_layout = centered_image_presentation.slide_layouts[6]

    centered_image_presentation_slides: Slides = centered_image_presentation.slides
    title_slide: Slide = centered_image_presentation_slides.add_slide(
        slide_layout)

    slide_width: int = centered_image_presentation.slide_width
    slide_height: int = centered_image_presentation.slide_height

    textbox_width = int(slide_width * 0.8)
    textbox_height = Inches(2)
    left = int((slide_width - textbox_width) / 2)

    title_shapes: SlideShapes = title_slide.shapes
    textbox: BaseShape = title_shapes.add_textbox(
        left, 0, textbox_width, textbox_height)
    text_frame: TextFrame = textbox.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title_text

    run.font.size = Pt(52)
    run.font.name = 'Arial'

    p.alignment = PP_ALIGN.CENTER

    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    for paragraph in text_frame.paragraphs:
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)

    text_height = Pt(52)
    new_top = int((slide_height - text_height) / 2)
    textbox.top = new_top
    textbox.left = int((slide_width - textbox_width) / 2)
    textbox.width = textbox_width
    textbox.height = textbox_height + text_height


def add_centered_image_slide(centered_image_presentation: presentation.Presentation, image_path: str) -> None:
    slide_layout = presentation.slide_layouts[6]
    centered_image_presentation_slides: Slides = centered_image_presentation.slides
    centered_image_slide: Slide = centered_image_presentation_slides.add_slide(
        slide_layout)

    img = Image.open(image_path)
    img_width_px, img_height_px = img.size
    img_aspect_ratio = img_width_px / img_height_px

    slide_width_emu = centered_image_presentation.slide_width
    slide_height_emu = centered_image_presentation.slide_height

    img_width_emu = int(slide_width_emu * 0.9)
    img_height_emu = int(img_width_emu / img_aspect_ratio)
    if img_height_emu > slide_height_emu * 0.9:
        img_height_emu = int(slide_height_emu * 0.9)
        img_width_emu = int(img_height_emu * img_aspect_ratio)

    left_emu = int((slide_width_emu - img_width_emu) / 2)
    top_emu = int((slide_height_emu - img_height_emu) / 2)

    centered_image_slide.shapes.add_picture(image_path, Emu(
        left_emu), Emu(top_emu), Emu(img_width_emu), Emu(img_height_emu))


def create_centered_images_presentation(directory: str, output_file: str) -> None:
    centered_image_presentation: presentation.Presentation = Presentation()
    centered_image_presentation.slide_width = Inches(10)
    centered_image_presentation.slide_height = Inches(5.625)
    processed_courses = set()

    for root, _, files in os.walk(directory):
        path_parts = root.split(os.sep)

        if len(path_parts) < 3 or path_parts[-1] == 'figures' or path_parts[-1] == 'presurvey':
            continue

        print(path_parts)
        course_id, course_run = path_parts[-2], path_parts[-1]
        course_key = (course_id, course_run)

        if course_key not in processed_courses:
            if course_id not in processed_courses:
                add_title_slide(centered_image_presentation, course_id)
                processed_courses.add(course_id)
            add_title_slide(centered_image_presentation, course_run)
            processed_courses.add(course_key)

        for file in sorted(files):
            if file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.wmf')):
                image_path = os.path.join(root, file)
                add_centered_image_slide(
                    centered_image_presentation, image_path)

    centered_image_presentation.save(output_file)
    print(f"""Presentation saved as {output_file} with {
          len(centered_image_presentation.slides)} slides.""")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Create a PowerPoint presentation from images in a directory, with titles for new course IDs and runs.")
    parser.add_argument('--directory', type=str, default=os.getcwd(),
                        help="The directory containing images.")
    parser.add_argument('--filename', type=str, default='presentation.pptx',
                        help="The filename for the output presentation.")
    args = parser.parse_args()

    directory: str = args.directory
    filename: str = args.filename

    create_centered_images_presentation(directory, filename)


if __name__ == "__main__":
    main()
